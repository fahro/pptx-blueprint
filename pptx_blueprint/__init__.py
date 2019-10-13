from pathlib import Path
import pptx
from typing import Union, Iterable, Tuple
from pptx.shapes.base import BaseShape
import subprocess
import os
_Pathlike = Union[str, Path]


class Template:
    """Helper class for modifying pptx templates.
    """

    def __init__(self, filename: _Pathlike) -> None:
        """Initializes a Template-Modifier.

        Args:
            filename (path-like): file name or path
        """
        self._template_path = filename
        self._presentation = pptx.Presentation(filename)
        pass

    def replace_text(self, label: str, new_text: str) -> None:
        """Replaces text placeholders on one or many slides.

        Args:
            label (str): label of the placeholder (without curly braces)
            text (str): new content
            scope: None, slide number, Slide object or iterable of Slide objects
        """
        slide_number, tag_name = self._parse_label(label)
        shapes = self._find_shapes(slide_number, tag_name)
        for shape in shapes:
            shape.text = new_text

    def replace_picture(self, label: str, filename: _Pathlike) -> None:
        """Replaces rectangle placeholders on one or many slides.

        Args:
            label (str): label of the placeholder (without curly braces)
            filename (path-like): path to an image file
        """
        pass

    def replace_table(self, label: str, data) -> None:
        """Replaces rectangle placeholders on one or many slides.

        Args:
            label (str): label of the placeholder (without curly braces)
            data (pandas.DataFrame): table to be inserted into the presentation
        """
        pass

    def _parse_label(self, label: str) -> Tuple[Union[int, str], str]:
        slide_number, tag_name = label.split(':')
        return int(slide_number) if slide_number != '*' else slide_number, tag_name

    def _find_shapes(self,
                     slide_number: Union[int, str],
                     tag_name: str) -> Iterable[BaseShape]:
        """Finds all shapes that match the label

        Args:
            label (str): label of the placeholder (without curly braces)
        """
        matched_shapes = []

        def _find_shapes_in_slide(slide):
            return filter(lambda shape: shape.text == f'{{{tag_name}}}', slide.shapes)

        if slide_number == '*':
            slides = self._presentation.slides
        else:
            # in label we are using 1 based indexing
            slide_index = slide_number - 1
            if slide_index < 0 or slide_index >= len(self._presentation.slides):
                raise IndexError(f"Can't find slide number {slide_number}.")
            slides = [self._presentation.slides[slide_index]]

        for slide in slides:
            matched_shapes.extend(_find_shapes_in_slide(slide))

        return matched_shapes

    def save(self, filename: _Pathlike) -> None:
        """Saves the updated pptx to the specified filepath.

         Args:
            filename (path-like): file name or path
        """
        # TODO: make sure that the user does not override the self._template_path
        pass

    def save_pdf(self, file_path: _Pathlike) -> None:
        """Exports the updated pptx to the specified filepath as pdf file.

        Args:
            filename (path-like) file name or path
        """

        try:

            subprocess.Popen(['libreoffice', '--version'],
                             stdout=subprocess.DEVNULL)  # check if libreoffice is installed

            path = Path(file_path)
            outdir = path.parent
            file_name = path.name

            # create temporary directory for pptx
            os.path.exists('tmp') or os.mkdir('tmp')
            template_temporary_path = f'tmp/{file_name}'
            # save current Template as pptx in temporary directory
            # TODO replace with self.save() method
            self._presentation.save(template_temporary_path)

            export_cmd = ['libreoffice', '--headless', '--convert-to',
                          'pdf', '--outdir', outdir, template_temporary_path]
            p = subprocess.Popen(
                export_cmd, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
            p.communicate()

            # remove temporary directory with pptx file
            os.remove(template_temporary_path)
            Path.rmdir(Path('tmp'))
        except FileNotFoundError:
            print("Libre Office not found.")
